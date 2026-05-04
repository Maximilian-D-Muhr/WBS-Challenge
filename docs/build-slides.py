"""ChallengeTracker presentation deck — Tracker Trio."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Theme ----
BG = RGBColor(0x0F, 0x14, 0x18)        # near-black
SURFACE = RGBColor(0x1A, 0x21, 0x27)   # card
ACCENT = RGBColor(0x4A, 0xDE, 0x80)    # green (success)
ACCENT2 = RGBColor(0x60, 0xA5, 0xFA)   # blue (info)
ACCENT3 = RGBColor(0xFB, 0xBF, 0x24)   # yellow (warning)
TEXT = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
CODE = RGBColor(0xCB, 0xD5, 0xE1)

# Member colors for name tag
COLORS = {
    "Rud": ACCENT2,
    "Umi": ACCENT3,
    "Max": ACCENT,
    "Tracker Trio": TEXT,
}

# ---- Setup ----
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK_LAYOUT = prs.slide_layouts[6]


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, font="Inter", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=TEXT,
                bullet_color=ACCENT, font="Inter", line_spacing=1.4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Bullet
        bullet = p.add_run()
        bullet.text = "▸  "
        bullet.font.name = font
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = bullet_color
        # Text
        run = p.add_run()
        run.text = item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_code_block(slide, x, y, w, h, lines, *, size=14):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = SURFACE
    # Adjust corner radius
    box.adjustments[0] = 0.08
    # Inner text
    pad = Emu(150000)
    tb = slide.shapes.add_textbox(x + pad, y + pad, w - 2 * pad, h - 2 * pad)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        run = p.add_run()
        run.text = line
        run.font.name = "JetBrains Mono"
        run.font.size = Pt(size)
        run.font.color.rgb = CODE


def add_name_tag(slide, name):
    """Bottom-right corner tag identifying which team member owns this slide."""
    color = COLORS.get(name, TEXT)
    # Pill background
    pill_w, pill_h = Inches(2.0), Inches(0.45)
    margin = Inches(0.4)
    pill_x = SW - pill_w - margin
    pill_y = SH - pill_h - margin
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_x, pill_y, pill_w, pill_h)
    pill.adjustments[0] = 0.5
    pill.line.color.rgb = color
    pill.line.width = Pt(1.5)
    pill.fill.solid()
    pill.fill.fore_color.rgb = SURFACE
    # Dot
    dot_size = Inches(0.18)
    dot_x = pill_x + Inches(0.18)
    dot_y = pill_y + (pill_h - dot_size) / 2
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, dot_y, dot_size, dot_size)
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    # Label
    label_x = pill_x + Inches(0.42)
    add_text(slide, label_x, pill_y, pill_w - Inches(0.5), pill_h, name,
             size=12, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)


def add_section_strip(slide, label, color):
    """Thin colored strip on the left edge identifying the section."""
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), SH)
    strip.line.fill.background()
    strip.fill.solid()
    strip.fill.fore_color.rgb = color


def add_slide_number(slide, n, total):
    add_text(slide, Inches(0.4), SH - Inches(0.6), Inches(2), Inches(0.4),
             f"{n:02d} / {total:02d}", size=11, color=MUTED)


def make_slide(owner, title_text, subtitle, content_fn, n, total):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    add_section_strip(slide, owner, COLORS[owner])
    # Title
    add_text(slide, Inches(0.7), Inches(0.7), Inches(11.5), Inches(0.7),
             title_text, size=36, color=TEXT, bold=True)
    # Subtitle / kicker
    if subtitle:
        add_text(slide, Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.4),
                 subtitle, size=16, color=COLORS[owner], bold=True)
    # Content callback
    content_fn(slide)
    add_name_tag(slide, owner)
    add_slide_number(slide, n, total)
    return slide


# ---- TITLE SLIDE ----
def title_slide():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    # Big logo dot
    dot_size = Inches(0.6)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.7), Inches(0.7), dot_size, dot_size)
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    add_text(slide, Inches(1.45), Inches(0.7), Inches(6), Inches(0.6),
             "ChallengeTracker", size=22, color=TEXT, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    # Big title
    add_text(slide, Inches(0.7), Inches(2.4), Inches(12), Inches(1.6),
             "Daily challenges, real progress,", size=54, color=TEXT, bold=True)
    add_text(slide, Inches(0.7), Inches(3.4), Inches(12), Inches(1.6),
             "live leaderboards.", size=54, color=ACCENT, bold=True)
    # Subline
    add_text(slide, Inches(0.7), Inches(5.0), Inches(12), Inches(0.6),
             "Fullstack group project · ASP.NET Core 10 · React 19 · two weeks",
             size=20, color=MUTED)
    # Team strip — three colored pills bottom-left
    pill_y = Inches(6.4)
    pill_h = Inches(0.5)
    pill_w = Inches(1.55)
    gap = Inches(0.15)
    members = [("Rud", ACCENT2), ("Umi", ACCENT3), ("Max", ACCENT)]
    x = Inches(0.7)
    for name, color in members:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, pill_y, pill_w, pill_h)
        pill.adjustments[0] = 0.5
        pill.line.color.rgb = color
        pill.line.width = Pt(1.5)
        pill.fill.solid()
        pill.fill.fore_color.rgb = SURFACE
        dot_s = Inches(0.18)
        dot_x = x + Inches(0.2)
        dot_y = pill_y + (pill_h - dot_s) / 2
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, dot_y, dot_s, dot_s)
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        add_text(slide, x + Inches(0.45), pill_y, pill_w - Inches(0.5), pill_h,
                 name, size=14, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        x += pill_w + gap
    # Tag right
    add_text(slide, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
             "Tracker Trio · WBS Coding School SE-004",
             size=12, color=MUTED)
    return slide


# Build deck
title_slide()

TOTAL = 11  # 1 title + 10 content


# ---- 02 · Rud — Backend Architecture ----
def s_rud_arch(slide):
    add_bullets(slide, Inches(0.7), Inches(2.1), Inches(7.5), Inches(4),
                [
                    "Three-layer Minimal API: Api / Application / Infrastructure",
                    "EF Core + SQLite, switchable to Azure SQL via connection string",
                    "ApplicationDbContext extends IdentityDbContext<User, IdentityRole<Guid>, Guid>",
                    "Single InitialCreate migration covers Identity + four domain tables",
                    "Scalar UI on /scalar/v1, ProblemDetails (RFC 7807) on every error",
                ], size=18, bullet_color=ACCENT2)
    add_code_block(slide, Inches(8.4), Inches(2.1), Inches(4.3), Inches(3.4),
                   [
                       "Api/",
                       "  Endpoints/",
                       "  Filters/",
                       "  Dtos/",
                       "Application/",
                       "  Interfaces/",
                       "  Services/",
                       "Infrastructure/",
                       "  ApplicationDbContext.cs",
                       "  Data/DbSeeder.cs",
                       "Models/",
                       "Migrations/",
                   ], size=12)


make_slide("Rud", "Backend Architecture", "Foundation · DbContext · Migrations",
           s_rud_arch, 2, TOTAL)


# ---- 03 · Umi — Domain Model & State Machine ----
def s_umi_domain(slide):
    add_bullets(slide, Inches(0.7), Inches(2.1), Inches(7), Inches(4),
                [
                    "Four entities: User, Challenge, Membership, ProgressEntry",
                    "Challenge owns memberships and progress entries",
                    "ChallengeStatus state machine: Open → Running → Completed",
                    "Owner-only transitions; ArgumentException maps to 400 ProblemDetails",
                    "Unique index (UserId, ChallengeId, LoggedAt) → 1-per-day rule",
                ], size=18, bullet_color=ACCENT3)
    # State machine visual
    sm_y = Inches(2.3)
    box_w, box_h = Inches(1.4), Inches(0.7)
    gap = Inches(0.25)
    x = Inches(8.0)
    for label, color in [("Open", MUTED), ("Running", ACCENT3), ("Completed", ACCENT)]:
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, sm_y, box_w, box_h)
        b.adjustments[0] = 0.3
        b.line.color.rgb = color
        b.line.width = Pt(1.5)
        b.fill.solid()
        b.fill.fore_color.rgb = SURFACE
        add_text(slide, x, sm_y, box_w, box_h, label, size=13, color=color,
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += box_w + gap
    # Arrows
    arr_y = sm_y + box_h / 2
    for i in range(2):
        ax = Inches(8.0) + (i + 1) * box_w + i * gap
        add_text(slide, ax - Inches(0.05), arr_y - Inches(0.18),
                 gap, Inches(0.4), "→", size=20, color=MUTED, bold=True,
                 align=PP_ALIGN.CENTER)
    # Code snippet
    add_code_block(slide, Inches(8.0), Inches(3.5), Inches(4.5), Inches(2.5),
                   [
                       "if (challenge.Status != ChallengeStatus.Open)",
                       "  throw new ArgumentException(",
                       '    "Challenge cannot be started.");',
                       "",
                       "challenge.Status = ChallengeStatus.Running;",
                       "await _db.SaveChangesAsync();",
                   ], size=12)


make_slide("Umi", "Domain Model & State Machine", "Entities · Status transitions · 1-per-day",
           s_umi_domain, 3, TOTAL)


# ---- 04 · Max — Frontend Stack ----
def s_max_frontend(slide):
    add_bullets(slide, Inches(0.7), Inches(2.1), Inches(7), Inches(4),
                [
                    "Vite + React 19 + TypeScript strict",
                    "Tailwind 4 + DaisyUI for instant design language",
                    "TanStack Query for fetching, caching, mutations",
                    "Zod for form schemas with cross-field rules",
                    "React Router v7 for routing & ProtectedRoute",
                    "react-toastify for success / error / 429 messages",
                ], size=18, bullet_color=ACCENT)
    # Component card
    card_x, card_y = Inches(8.0), Inches(2.1)
    card_w, card_h = Inches(4.7), Inches(4.2)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  card_x, card_y, card_w, card_h)
    card.adjustments[0] = 0.05
    card.line.color.rgb = ACCENT
    card.line.width = Pt(1)
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE
    add_text(slide, card_x + Inches(0.3), card_y + Inches(0.25),
             card_w - Inches(0.6), Inches(0.5),
             "frontend/src/", size=14, color=ACCENT, bold=True)
    add_bullets(slide, card_x + Inches(0.3), card_y + Inches(0.85),
                card_w - Inches(0.6), Inches(3.2),
                [
                    "components/  Reusable UI",
                    "pages/  Route components",
                    "data/  TanStack Query hooks",
                    "contexts/  AuthContext",
                    "utils/  api.ts, problemDetails.ts",
                    "types/  Mirrors backend models",
                ], size=14, bullet_color=ACCENT, line_spacing=1.5)


make_slide("Max", "Frontend Stack", "Vite · React 19 · TanStack Query · Zod",
           s_max_frontend, 4, TOTAL)


# ---- 05 · Rud — Auth & JWT ----
def s_rud_auth(slide):
    add_bullets(slide, Inches(0.7), Inches(2.1), Inches(7.5), Inches(4),
                [
                    "ASP.NET Core Identity for password hashing & user store",
                    "JWT Bearer with 64-char symmetric key from configuration",
                    "Token claims: sub (user id), email, displayName",
                    "POST /auth/register, POST /auth/login, GET /auth/me",
                    "Custom WithValidation<T>() filter → 400 ProblemDetails on invalid DTOs",
                ], size=18, bullet_color=ACCENT2)
    add_code_block(slide, Inches(8.4), Inches(2.1), Inches(4.3), Inches(3.6),
                   [
                       "var claims = new[]",
                       "{",
                       "  new Claim(",
                       "    JwtRegisteredClaimNames.Sub,",
                       "    user.Id.ToString()),",
                       "  new Claim(",
                       "    JwtRegisteredClaimNames.Email,",
                       "    user.Email),",
                       "  new Claim(\"displayName\",",
                       "    user.DisplayName)",
                       "};",
                   ], size=12)


make_slide("Rud", "Authentication & JWT", "Identity · Bearer tokens · Validation filter",
           s_rud_auth, 5, TOTAL)


# ---- 06 · Umi — Memberships & Progress ----
def s_umi_progress(slide):
    add_bullets(slide, Inches(0.7), Inches(2.1), Inches(7.5), Inches(4),
                [
                    "Public challenge → Active immediately; private → Pending until owner approves",
                    "Re-join after Leave reactivates the existing row (avoids unique-index conflict)",
                    "Progress logging stack: challenge Running + active membership + within window",
                    "1-per-day pre-check returns clean 409, unique index is the safety net",
                    "Edit & delete window: owner-only, 24h after creation",
                ], size=18, bullet_color=ACCENT3)
    add_code_block(slide, Inches(8.4), Inches(2.1), Inches(4.3), Inches(3.6),
                   [
                       "var existing = await _db",
                       "  .ProgressEntries.AnyAsync(p =>",
                       "    p.UserId == userId &&",
                       "    p.ChallengeId == dto.ChallengeId &&",
                       "    p.LoggedAt == loggedAt);",
                       "",
                       "if (existing)",
                       "  throw new InvalidOpException(",
                       '    "Already logged today.");',
                   ], size=12)


make_slide("Umi", "Memberships & Progress", "Join/Leave · 1-per-day · 24h edit window",
           s_umi_progress, 6, TOTAL)


# ---- 07 · Max — UI Components ----
def s_max_components(slide):
    add_bullets(slide, Inches(0.7), Inches(2.1), Inches(7.5), Inches(4),
                [
                    "ChallengeCard, ChallengeStateBadge — reusable visual building blocks",
                    "JoinLeaveButton flips action based on useMyMembership query",
                    "OwnerActions renders Start when Open, Complete when Running",
                    "LogProgressForm with Zod validation + 409/429-aware toasts",
                    "Leaderboard table with rank, avatar initials, totals",
                    "ErrorBoundary + NotFoundPage for graceful failure",
                ], size=17, bullet_color=ACCENT)
    add_code_block(slide, Inches(8.4), Inches(2.1), Inches(4.3), Inches(3.6),
                   [
                       "function useLogProgress(",
                       "  challengeId: string) {",
                       "  return useMutation({",
                       "    mutationFn: payload =>",
                       "      api(`/progress-entries`,",
                       "        { method: \"POST\",",
                       "          body: payload }),",
                       "    onSuccess: () => qc",
                       "      .invalidateQueries(...)",
                       "  });",
                       "}",
                   ], size=12)


make_slide("Max", "UI Components & Mutations", "TanStack Query · Optimistic UX · Toasts",
           s_max_components, 7, TOTAL)


# ---- 08 · Rud — Docker + Health + CORS ----
def s_rud_docker(slide):
    add_bullets(slide, Inches(0.7), Inches(2.1), Inches(7), Inches(4),
                [
                    "Multi-stage Dockerfile: SDK image builds, slim aspnet runtime hosts",
                    "docker-compose.yml runs the API on :8080 with named SQLite volume",
                    "Healthcheck pings /health every 15s — Compose reports container state",
                    "/health = liveness, /ready = DB connectivity + pending migrations",
                    "CORS policy reads origins from env so deployment can extend it",
                    "Both endpoints exempt from rate limiting and authentication",
                ], size=17, bullet_color=ACCENT2)
    add_code_block(slide, Inches(8.0), Inches(2.1), Inches(4.7), Inches(4),
                   [
                       "services:",
                       "  api:",
                       "    build: ./backend/...",
                       "    ports: [\"8080:8080\"]",
                       "    environment:",
                       "      Jwt__Key: ${JWT_KEY}",
                       "      ConnectionStrings__",
                       "        DefaultConnection: ...",
                       "    volumes:",
                       "      - challenge-data:/app/data",
                       "    healthcheck:",
                       "      test: [/health]",
                   ], size=11)


make_slide("Rud", "Docker · Health · CORS", "One-command stack · Production-ready",
           s_rud_docker, 8, TOTAL)


# ---- 09 · Umi — Leaderboard + Rate Limiting ----
def s_umi_rate(slide):
    add_bullets(slide, Inches(0.7), Inches(2.1), Inches(7), Inches(4),
                [
                    "Leaderboard groups progress entries by user, sums amounts, ranks",
                    "Public endpoint (no auth) — embeddable on a future landing page",
                    "Two rate-limit policies: progress-post 10/min, general-post 5/min",
                    "Partition key = NameIdentifier claim → per-user buckets",
                    "UseRateLimiter() runs after UseAuthentication so claims resolve",
                    "Rejection writes 429 + Retry-After + ProblemDetails body",
                ], size=17, bullet_color=ACCENT3)
    add_code_block(slide, Inches(8.0), Inches(2.1), Inches(4.7), Inches(4.2),
                   [
                       "options.AddPolicy(",
                       "  \"progress-post\", ctx =>",
                       "  RateLimitPartition",
                       "    .GetFixedWindowLimiter(",
                       "      partitionKey: GetUserId(ctx),",
                       "      factory: _ => new ()",
                       "      {",
                       "        PermitLimit = 10,",
                       "        Window = TimeSpan",
                       "          .FromMinutes(1)",
                       "      }));",
                   ], size=11)


make_slide("Umi", "Leaderboard & Rate Limiting", "Aggregate queries · Per-user 429s",
           s_umi_rate, 9, TOTAL)


# ---- 10 · Max — Live Demo & Closing ----
def s_max_demo(slide):
    add_bullets(slide, Inches(0.7), Inches(2.1), Inches(7.5), Inches(4),
                [
                    "Login as Alice → start the seeded challenge",
                    "Log in as Bob → see Leave button (active member detection works)",
                    "Log progress 25 pages → toast confirms, leaderboard updates",
                    "Log again same day → 409 toast \"already logged for today\"",
                    "Spam the button → 429 toast with Retry-After hint",
                    "Visit /not-a-route → NotFoundPage rendered inside the layout",
                ], size=17, bullet_color=ACCENT)
    # Demo URL card
    card_x, card_y = Inches(8.4), Inches(2.1)
    card_w, card_h = Inches(4.3), Inches(3.4)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  card_x, card_y, card_w, card_h)
    card.adjustments[0] = 0.06
    card.line.color.rgb = ACCENT
    card.line.width = Pt(1)
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE
    add_text(slide, card_x + Inches(0.3), card_y + Inches(0.25),
             card_w - Inches(0.6), Inches(0.45),
             "LIVE DEMO", size=12, color=ACCENT, bold=True)
    add_text(slide, card_x + Inches(0.3), card_y + Inches(0.75),
             card_w - Inches(0.6), Inches(0.5),
             "localhost:5173", size=22, color=TEXT, bold=True,
             font="JetBrains Mono")
    add_text(slide, card_x + Inches(0.3), card_y + Inches(1.4),
             card_w - Inches(0.6), Inches(0.4),
             "Seed users", size=12, color=MUTED, bold=True)
    add_text(slide, card_x + Inches(0.3), card_y + Inches(1.75),
             card_w - Inches(0.6), Inches(0.4),
             "alice@example.com", size=14, color=TEXT,
             font="JetBrains Mono")
    add_text(slide, card_x + Inches(0.3), card_y + Inches(2.15),
             card_w - Inches(0.6), Inches(0.4),
             "bob@example.com", size=14, color=TEXT,
             font="JetBrains Mono")
    add_text(slide, card_x + Inches(0.3), card_y + Inches(2.6),
             card_w - Inches(0.6), Inches(0.4),
             "Password1!", size=14, color=ACCENT,
             font="JetBrains Mono", bold=True)


make_slide("Max", "Live Demo", "End-to-end flow · localhost:5173",
           s_max_demo, 10, TOTAL)


# ---- 11 · Closing — Tracker Trio ----
def s_closing(slide):
    # Override the default by not using the section strip — but we need a closing
    pass

slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide)
add_text(slide, Inches(0.7), Inches(0.7), Inches(11.5), Inches(0.7),
         "Thanks", size=36, color=TEXT, bold=True)
add_text(slide, Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.4),
         "Tracker Trio · WBS Coding School SE-004", size=16, color=ACCENT, bold=True)

# What's next box
box_x, box_y = Inches(0.7), Inches(2.4)
box_w, box_h = Inches(5.8), Inches(4.2)
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, box_y, box_w, box_h)
box.adjustments[0] = 0.04
box.line.color.rgb = ACCENT2
box.line.width = Pt(1)
box.fill.solid()
box.fill.fore_color.rgb = SURFACE
add_text(slide, box_x + Inches(0.4), box_y + Inches(0.3), box_w - Inches(0.8), Inches(0.5),
         "What's next", size=18, color=ACCENT2, bold=True)
add_bullets(slide, box_x + Inches(0.4), box_y + Inches(0.95),
            box_w - Inches(0.8), Inches(3.2),
            [
                "Optimistic cache surgery (zero-flicker mutations)",
                "Daily / weekly leaderboard periods",
                "Astro landing page on the root domain",
                "Azure deployment: App Service + SQL + SWA",
            ], size=15, bullet_color=ACCENT2)

# Repo box
box2_x = Inches(7.0)
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box2_x, box_y, box_w, box_h)
box2.adjustments[0] = 0.04
box2.line.color.rgb = ACCENT
box2.line.width = Pt(1)
box2.fill.solid()
box2.fill.fore_color.rgb = SURFACE
add_text(slide, box2_x + Inches(0.4), box_y + Inches(0.3),
         box_w - Inches(0.8), Inches(0.5),
         "Repository", size=18, color=ACCENT, bold=True)
add_text(slide, box2_x + Inches(0.4), box_y + Inches(0.95),
         box_w - Inches(0.8), Inches(0.5),
         "github.com/Maximilian-D-Muhr/", size=14, color=TEXT,
         font="JetBrains Mono")
add_text(slide, box2_x + Inches(0.4), box_y + Inches(1.3),
         box_w - Inches(0.8), Inches(0.5),
         "  WBS-Challenge", size=14, color=TEXT,
         font="JetBrains Mono")
add_text(slide, box2_x + Inches(0.4), box_y + Inches(2.0),
         box_w - Inches(0.8), Inches(0.5),
         "Quick start", size=14, color=ACCENT, bold=True)
add_text(slide, box2_x + Inches(0.4), box_y + Inches(2.4),
         box_w - Inches(0.8), Inches(0.5),
         "$ docker compose up", size=13, color=CODE,
         font="JetBrains Mono")
add_text(slide, box2_x + Inches(0.4), box_y + Inches(2.8),
         box_w - Inches(0.8), Inches(0.5),
         "$ npm run dev", size=13, color=CODE,
         font="JetBrains Mono")

# Three member tags bottom
pill_y = Inches(6.8)
pill_h = Inches(0.45)
pill_w = Inches(1.55)
gap = Inches(0.15)
members = [("Rud", ACCENT2), ("Umi", ACCENT3), ("Max", ACCENT)]
total_w = pill_w * len(members) + gap * (len(members) - 1)
x = (SW - total_w) / 2
for name, color in members:
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, pill_y, pill_w, pill_h)
    pill.adjustments[0] = 0.5
    pill.line.color.rgb = color
    pill.line.width = Pt(1.5)
    pill.fill.solid()
    pill.fill.fore_color.rgb = SURFACE
    dot_s = Inches(0.18)
    dot_x = x + Inches(0.18)
    dot_y = pill_y + (pill_h - dot_s) / 2
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, dot_y, dot_s, dot_s)
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    add_text(slide, x + Inches(0.42), pill_y, pill_w - Inches(0.5), pill_h,
             name, size=12, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    x += pill_w + gap

add_slide_number(slide, 11, TOTAL)


# Save
import os
os.makedirs("/tmp/ct-slides", exist_ok=True)
out = "/tmp/ct-slides/ChallengeTracker-Demo.pptx"
prs.save(out)
print(f"Saved: {out}")
